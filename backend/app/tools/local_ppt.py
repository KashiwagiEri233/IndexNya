"""本地 PPT 生成器 — 不调用模型和在线 PPT 服务。"""
from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT_DIR = Path(__file__).resolve().parent.parent.parent / "outputs" / "ppts"

NAVY = RGBColor(38, 58, 70)
TEAL = RGBColor(81, 167, 160)
MINT = RGBColor(229, 245, 241)
CORAL = RGBColor(238, 151, 130)
INK = RGBColor(47, 65, 74)
MUTED = RGBColor(111, 129, 136)
WHITE = RGBColor(255, 255, 255)


def _add_textbox(slide, left, top, width, height, text="", *, size=22, color=INK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def _add_header(slide, title: str, number: int):
    _add_textbox(slide, Inches(0.65), Inches(0.38), Inches(11.5), Inches(0.55), title, size=25, color=NAVY, bold=True)
    line = slide.shapes.add_shape(1, Inches(0.7), Inches(1.1), Inches(11.9), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = MINT
    line.line.fill.background()
    _add_textbox(slide, Inches(11.8), Inches(7.05), Inches(0.7), Inches(0.25), str(number), size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def _add_bullets(slide, items: list[str], *, top=Inches(1.45), accent=TEAL):
    box = slide.shapes.add_textbox(Inches(0.85), top, Inches(11.4), Inches(5.1))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.1)
    frame.margin_right = Inches(0.1)
    frame.margin_top = Inches(0.08)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.level = 0
        paragraph.space_after = Pt(14)
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = INK
    # 左侧色条
    bar = slide.shapes.add_shape(1, Inches(0.65), top, Inches(0.08), Inches(4.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()


def _add_callout(slide, text: str, top=Inches(5.65)):
    shape = slide.shapes.add_shape(5, Inches(0.85), top, Inches(11.4), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MINT
    shape.line.color.rgb = MINT
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = NAVY


def _build_outline(topic: str, context: list[str]) -> list[tuple[str, list[str], str | None]]:
    context_items = context[:4] or [f"本次课件围绕“{topic}”展开。"]
    return [
        ("学习目标", [f"理解“{topic}”的基本定义", "掌握核心概念与相互关系", "能够用一个简单例子说明它的作用"], "先建立整体认识，再进入细节。"),
        ("本次学习背景", context_items, "课件内容根据当前对话和学习记录整理。"),
        ("核心概念", [f"{topic}解决什么问题", "它由哪些关键部分组成", "这些部分之间如何配合"], "把复杂内容拆成几个可以分别理解的小块。"),
        ("学习要点", [f"定义：用一句话说明{topic}", "过程：按顺序列出关键步骤", "判断：如何看出结果是否合理"], "遇到新概念时，可以按照“定义—过程—判断”来梳理。"),
        ("示例与应用", [f"从一个常见问题出发理解{topic}", "按照步骤完成一次简单演示", "观察结果并说明原因"], "先看一个小例子，再迁移到更复杂的问题。"),
        ("常见误区", ["只记结论，不理解适用条件", "跳过中间步骤，导致理解断层", "忽略输入、限制条件和结果检查"], "学习时把“为什么”与“什么时候不能用”一起记下来。"),
        ("小结", [f"{topic}的核心是解决一个具体问题", "掌握定义、步骤、例子和边界", "下一步可以通过练习或资料继续巩固"], "把今天的内容整理成自己的笔记。"),
    ]


def _context_items(profile: dict | None, extra: str) -> list[str]:
    items: list[str] = []
    for key, label in (("major", "专业方向"), ("learning_goals", "学习目标"), ("knowledge_base", "基础情况"), ("common_mistakes", "常见难点")):
        value = (profile or {}).get(key)
        if value:
            items.append(f"{label}：{str(value)[:100]}")
    history = extra.split("对话历史：", 1)[-1] if "对话历史：" in extra else ""
    history = history.split("额外要求：", 1)[0]
    for line in history.splitlines():
        line = line.strip()
        if line.startswith("user:") or line.startswith("用户:"):
            text = line.split(":", 1)[-1].strip()
            if text and text not in items:
                items.append(f"对话中提出的问题：{text[:120]}")
    return items[:4]


def generate_local_ppt(topic: str, *, profile: dict | None = None, extra: str = "") -> dict[str, Any]:
    """使用固定教学模板生成本地 .pptx 文件，并带入当前画像和对话上下文。"""
    topic = (topic or "学习主题").strip()[:80]
    context = _context_items(profile, extra)
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 标题页
    slide = prs.slides.add_slide(blank)
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = MINT
    _add_textbox(slide, Inches(0.9), Inches(1.4), Inches(11.5), Inches(1.1), topic, size=38, color=NAVY, bold=True)
    _add_textbox(slide, Inches(0.95), Inches(2.7), Inches(10.5), Inches(0.55), "Index 学习岛 · 本地生成课件", size=20, color=TEAL, bold=True)
    _add_textbox(slide, Inches(0.95), Inches(5.85), Inches(10.5), Inches(0.5), datetime.now().strftime("%Y年%m月%d日"), size=14, color=MUTED)
    circle = slide.shapes.add_shape(9, Inches(10.4), Inches(0.8), Inches(2.0), Inches(2.0))
    circle.fill.solid(); circle.fill.fore_color.rgb = CORAL; circle.line.fill.background()

    for number, (title, bullets, callout) in enumerate(_build_outline(topic, context), start=2):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE
        _add_header(slide, title, number)
        _add_bullets(slide, bullets, accent=CORAL if number % 2 == 0 else TEAL)
        _add_callout(slide, callout)

    filename = f"index-learning-{datetime.now().strftime('%Y%m%d-%H%M%S-%f')}.pptx"
    path = OUTPUT_DIR / filename
    prs.save(path)
    return {
        "topic": topic,
        "context_items": context,
        "ppt_path": str(path),
        "filename": filename,
        "status": "completed",
        "slide_count": len(prs.slides),
        "generator": "local-template",
    }
